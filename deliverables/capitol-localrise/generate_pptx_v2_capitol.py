"""
generate_pptx_v2_capitol.py  — REDESIGNED v2
Premium executive presentation | Capitol Consultoria x LocalRise Advisory
Design level: McKinsey / BCG / Bain
Author: LocalRise Advisory

Design System:
  Grid      : 1.8cm margin L/R | header 0-3.7cm | content 3.9-17.4cm | footer 17.4-19.05cm
  Colors    : BLACK #0A0A0A | CARD #141414 | RED #E31B23 | DRED #B51218
              WHITE #FFFFFF | SILVER #C4C6C8 | MUTE #686B6F | BORDER #252525
  Typography: Hero=52pt | Title=34pt | Head2=20pt | Head3=15pt | Body=12pt | Caption=9pt
  Logo      : Fixed in footer right (all interior slides) + cover special placement
"""
from pathlib import Path
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

# ─── Paths ──────────────────────────────────────────────────────────────────
BASE_DIR  = Path(__file__).resolve().parent
LOGO_PRI  = Path(r"C:\Users\digui\Documents\localrise-brain\Claude\Claude Local RIse\assents\logo Local Rise Principal.png")
LOGO_SEC  = Path(r"C:\Users\digui\Documents\localrise-brain\Claude\Claude Local RIse\assents\logo Local Rise Secundária.png")
LOGO_DARK = Path(r"C:\Users\digui\Documents\localrise-brain\Claude\Claude Local RIse\assents\Oficial Preta PNG (2).png")
OUT_PATH  = BASE_DIR / "apresentacao-capitol-localrise-v2.pptx"

# ─── Design Tokens ──────────────────────────────────────────────────────────
BLACK   = RGBColor(0x0A, 0x0A, 0x0A)   # slide background
CARD    = RGBColor(0x14, 0x14, 0x14)   # primary card surface
CARD2   = RGBColor(0x1E, 0x1E, 0x1E)   # elevated card
BORDER  = RGBColor(0x25, 0x25, 0x25)   # card border / dividers
RED     = RGBColor(0xE3, 0x1B, 0x23)   # primary accent
DRED    = RGBColor(0xB5, 0x12, 0x18)   # deep red (headers, depth)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)   # primary text
SILVER  = RGBColor(0xC4, 0xC6, 0xC8)  # secondary text
MUTE    = RGBColor(0x68, 0x6B, 0x6F)   # captions / tertiary
LGRAY   = RGBColor(0x3C, 0x3F, 0x42)   # decorative dividers / rules

# ─── Slide Canvas ───────────────────────────────────────────────────────────
W = 33.867   # cm — 16:9 FullHD
H = 19.05    # cm

# ─── Grid Constants ─────────────────────────────────────────────────────────
ML    = 1.8         # margin left
MR    = 1.8         # margin right
CW    = W - ML - MR # 30.267cm — usable content width
HDR_T = 0.5         # header title y
HDR_S = 2.1         # header subtitle y
HDR_L = 3.45        # header separator line y
CTT   = 3.9         # content top y
CTB   = 17.3        # content bottom y
CTH   = CTB - CTT   # 13.4cm — content height
FTR   = 17.45       # footer y
FTR_H = 19.05 - FTR # 1.6cm — footer height

# ─── Metadata ───────────────────────────────────────────────────────────────
DATE    = "08 de abril de 2026"
COMPANY = "Capitol Consultoria"


# ════════════════════════════════════════════════════════════════════════════
# PRIMITIVES
# ════════════════════════════════════════════════════════════════════════════

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = Cm(W)
    prs.slide_height = Cm(H)
    return prs


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color=BLACK):
    """Fill entire slide background."""
    s = slide.shapes.add_shape(1, Cm(0), Cm(0), Cm(W), Cm(H))
    s.line.fill.background(); s.line.width = 0
    s.fill.solid(); s.fill.fore_color.rgb = color
    return s


def box(slide, x, y, w, h, fill: RGBColor, border_color=None, border_pt=0.75):
    """Solid rectangle, optional 1-side border via thin rect trick."""
    s = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    s.line.fill.background(); s.line.width = 0
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border_color:
        from pptx.util import Pt as PPt
        s.line.color.rgb = border_color
        s.line.width = PPt(border_pt)
    return s


def rule(slide, x, y, w, color=RED, t=0.06):
    """Thin horizontal rule."""
    return box(slide, x, y, w, t, color)


def vline(slide, x, y, h, color=RED, t=0.06):
    """Thin vertical rule."""
    s = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(t), Cm(h))
    s.line.fill.background(); s.line.width = 0
    s.fill.solid(); s.fill.fore_color.rgb = color
    return s


def logo(slide, x, y, w=3.8, path=None):
    """Add LocalRise logo. path defaults to principal version."""
    p = path or LOGO_PRI
    if not p.exists():
        p = LOGO_PRI  # fallback
    if p.exists():
        slide.shapes.add_picture(str(p), Cm(x), Cm(y), width=Cm(w))


def tx(slide, text, x, y, w, h,
       size=13, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False,
       font="Calibri"):
    """Single-paragraph text box."""
    b = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    b.word_wrap = True
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return b


def txm(slide, lines, x, y, w, h,
        size=13, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, spacing_pt=None,
        font="Calibri"):
    """Multi-paragraph text box. lines = list of str."""
    b = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    b.word_wrap = True
    tf = b.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return b


# ════════════════════════════════════════════════════════════════════════════
# LAYOUT COMPONENTS
# ════════════════════════════════════════════════════════════════════════════

def slide_background(slide):
    """Full black background for all slides."""
    bg(slide, BLACK)


def slide_footer(slide, page=None):
    """
    Consistent footer: separator line + left caption + right logo.
    Applied to every interior slide.
    """
    # Separator line
    rule(slide, 0, FTR, W, LGRAY, t=0.05)

    # Left caption
    note = "Confidencial  |  LocalRise Advisory x Capitol Consultoria  |  Abril 2026"
    tx(slide, note, ML, FTR + 0.25, 22, 1.0,
       size=8.5, color=MUTE, align=PP_ALIGN.LEFT)

    # Page number (optional)
    if page:
        tx(slide, str(page), W/2 - 0.5, FTR + 0.25, 1.0, 1.0,
           size=8.5, color=MUTE, align=PP_ALIGN.CENTER)

    # Logo — right side of footer
    logo(slide, x=W - 4.4, y=FTR + 0.18, w=3.8)


def slide_header(slide, title, subtitle=None):
    """
    Interior slide header: red top bar + title + optional subtitle + separator.
    """
    # Top accent bar (thin red strip)
    box(slide, 0, 0, W, 0.22, RED)

    # Slide title
    tx(slide, title, ML, HDR_T, 28, 1.9,
       size=30, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Subtitle
    if subtitle:
        tx(slide, subtitle, ML, HDR_S, CW, 1.1,
           size=13, bold=False, color=SILVER, align=PP_ALIGN.LEFT)

    # Separator
    rule(slide, ML, HDR_L, CW, LGRAY, t=0.05)


def interior_slide(prs, title, subtitle=None, page=None):
    """
    Factory: creates a black slide with header + footer pre-applied.
    Returns the slide for content to be added to content zone.
    """
    sl = blank(prs)
    slide_background(sl)
    slide_header(sl, title, subtitle)
    slide_footer(sl, page)
    return sl


# ════════════════════════════════════════════════════════════════════════════
# CARD COMPONENTS
# ════════════════════════════════════════════════════════════════════════════

def card_full_row(slide, x, y, w, h, label=None, value=None,
                  desc=None, accent=RED):
    """Full-width card with left accent bar + label + value + description."""
    box(slide, x, y, w, h, CARD)
    vline(slide, x, y, h, accent, t=0.45)  # left accent
    if label:
        tx(slide, label, x + 0.75, y + 0.25, w - 1.0, 0.9,
           size=10.5, bold=True, color=RED)
    if value:
        tx(slide, value, x + 0.75, y + 1.1, 8, h - 1.3,
           size=22, bold=True, color=WHITE)
    if desc:
        tx(slide, desc, x + 9.0, y + (h / 2) - 0.4, w - 10.0, 1.5,
           size=12, color=SILVER)


def kpi_card(slide, x, y, w, h, number, label, sub=None):
    """KPI card: big number + label + optional sublabel."""
    box(slide, x, y, w, h, CARD)
    rule(slide, x, y, w, RED, t=0.25)  # top accent
    tx(slide, number, x + 0.4, y + 0.6, w - 0.8, 2.8,
       size=40, bold=True, color=RED, align=PP_ALIGN.CENTER)
    rule(slide, x + 1.0, y + 3.4, w - 2.0, LGRAY, t=0.05)
    tx(slide, label, x + 0.3, y + 3.6, w - 0.6, 1.6,
       size=11, color=SILVER, align=PP_ALIGN.CENTER)
    if sub:
        tx(slide, sub, x + 0.3, y + 4.8, w - 0.6, 1.0,
           size=9.5, color=MUTE, align=PP_ALIGN.CENTER)


def numbered_card(slide, x, y, w, h, num, title, body=None):
    """Card with large number badge + title + body."""
    box(slide, x, y, w, h, CARD)
    rule(slide, x, y, w, RED, t=0.25)  # top accent
    # Number
    tx(slide, num,
       x + 0.35, y + 0.45, 1.8, 1.5,
       size=28, bold=True, color=DRED, align=PP_ALIGN.LEFT)
    # Title
    tx(slide, title, x + 0.35, y + 2.0, w - 0.7, 1.4,
       size=15, bold=True, color=WHITE)
    # Body
    if body:
        rule(slide, x + 0.35, y + 3.5, w - 0.7, LGRAY, t=0.04)
        tx(slide, body, x + 0.35, y + 3.75, w - 0.7, h - 4.2,
           size=11.5, color=SILVER)


def phase_card(slide, x, y, w, h, num, title, deadline, items):
    """Phase card: number + title + deadline badge + bullet items."""
    box(slide, x, y, w, h, CARD)
    rule(slide, x, y, w, RED, t=0.25)
    # Phase number
    tx(slide, num, x + 0.4, y + 0.5, 5, 1.2,
       size=11, bold=True, color=RED)
    # Deadline badge
    box(slide, x + w - 3.2, y + 0.45, 2.9, 1.0, CARD2)
    tx(slide, deadline, x + w - 3.1, y + 0.55, 2.8, 0.9,
       size=10, color=SILVER, align=PP_ALIGN.CENTER)
    # Title
    tx(slide, title, x + 0.4, y + 1.7, w - 0.8, 1.3,
       size=15, bold=True, color=WHITE)
    rule(slide, x + 0.4, y + 3.1, w - 0.8, LGRAY, t=0.04)
    # Items
    tx(slide, items, x + 0.4, y + 3.35, w - 0.8, h - 3.8,
       size=11.5, color=SILVER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDES
# ════════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    sl = blank(prs)

    # ── Backgrounds ──────────────────────────────────────────────────
    bg(sl, BLACK)

    # Right red panel (40% of width)
    red_x = W * 0.60
    box(sl, red_x, 0, W - red_x, H, DRED)
    # Thin transition strip
    box(sl, red_x - 0.15, 0, 0.35, H, RED)

    # ── Left content ─────────────────────────────────────────────────
    # Eyebrow
    tx(sl, "PROPOSTA ESTRATEGICA  |  " + DATE,
       2.0, 1.8, 18, 0.9,
       size=10, bold=False, color=MUTE)

    # Main headline — 2 lines
    tx(sl, "INTELIGENCIA",
       2.0, 2.8, 20, 3.5,
       size=54, bold=True, color=WHITE)
    tx(sl, "REGULATORIA",
       2.0, 5.6, 20, 3.5,
       size=54, bold=True, color=RED)

    rule(sl, 2.0, 9.25, 15.0, RED, t=0.12)

    tx(sl, "COMO VANTAGEM COMPETITIVA",
       2.0, 9.65, 18, 1.4,
       size=17, bold=False, color=SILVER)

    tx(sl, "Apresentado para " + COMPANY,
       2.0, 11.3, 18, 1.0,
       size=13, bold=False, color=MUTE)

    # Logo — bottom left, inside black panel
    logo(sl, x=2.0, y=15.8, w=4.5)

    # Presenter block
    tx(sl, "Nicolas Slater",
       2.0, 13.4, 14, 0.9,
       size=12, bold=True, color=WHITE)
    tx(sl, "Diretoria Comercial  |  LocalRise Advisory",
       2.0, 14.1, 14, 0.8,
       size=11, color=MUTE)

    # ── Right panel content ───────────────────────────────────────────
    px = red_x + 0.8
    pw = W - red_x - 1.2

    items = [
        ("01", "Presenca Digital"),
        ("02", "Automacao com IA"),
        ("03", "SaaS B2B"),
        ("04", "ROI em 90 dias"),
    ]
    for i, (n, label) in enumerate(items):
        iy = 3.5 + i * 3.4
        tx(sl, n, px, iy, 1.4, 1.2,
           size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        rule(sl, px, iy + 1.3, pw - 0.4, RGBColor(0xFF, 0xFF, 0xFF), t=0.04)
        tx(sl, label, px, iy + 1.6, pw - 0.2, 1.2,
           size=13, bold=True, color=WHITE)

    return sl


def slide_02_sumario(prs):
    sl = interior_slide(prs, "Sumario Executivo",
                        "O que vamos apresentar hoje", page=2)

    topics = [
        ("01", "Diagnostico",         "A situacao atual da Capitol"),
        ("02", "A Oportunidade",       "O que esta sendo deixado na mesa"),
        ("03", "A Solucao",            "Automacao + IA + presenca digital"),
        ("04", "Produtos SaaS",        "Portfolio de inteligencia regulatoria"),
        ("05", "Pipeline de IA",       "6 agentes, 1 sistema integrado"),
        ("06", "Implementacao",        "4 fases, 6 meses"),
        ("07", "ROI Esperado",         "Metricas e retorno projetado"),
        ("08", "Proposta Comercial",   "Investimento e proximos passos"),
    ]

    for i, (num, title, desc) in enumerate(topics):
        col = i % 2
        row = i // 2
        cx = ML + col * 15.3
        cy = CTT + row * 3.4

        # Card
        box(sl, cx, cy, 14.6, 3.0, CARD)
        vline(sl, cx, cy, 3.0, RED if col == 0 else DRED, t=0.4)

        tx(sl, num, cx + 0.7, cy + 0.4, 2.0, 1.0,
           size=22, bold=True, color=DRED)
        tx(sl, title, cx + 0.7, cy + 1.5, 13, 0.9,
           size=13, bold=True, color=WHITE)
        tx(sl, desc, cx + 0.7, cy + 2.3, 13, 0.8,
           size=10.5, color=MUTE)

    return sl


def slide_03_diagnostico(prs):
    sl = interior_slide(prs, "Diagnostico",
                        "Onde a Capitol Consultoria esta hoje", page=3)

    problems = [
        ("Presenca Digital\nMinima",
         "87 seguidores no Instagram.\nSem site. Sem blog.\nSem funil de captacao ativo."),
        ("Dependencia de\nIndicacoes",
         "100% dos clientes via networking.\nNenhum canal proprio.\nRisco alto de sazonalidade."),
        ("Operacao\nManual",
         "Monitoramento legislativo\nfeito a mao. Sem alertas.\nHoras perdidas toda semana."),
    ]

    for i, (title, body) in enumerate(problems):
        cx = ML + i * 10.1
        cy = CTT

        box(sl, cx, cy, 9.5, CTH, CARD)
        rule(sl, cx, cy, 9.5, RED, t=0.25)

        # Large problem number
        tx(sl, f"0{i+1}", cx + 0.5, cy + 0.55, 2.5, 2.0,
           size=40, bold=True, color=DRED)

        tx(sl, title, cx + 0.5, cy + 2.7, 8.5, 2.5,
           size=17, bold=True, color=WHITE)

        rule(sl, cx + 0.5, cy + 5.4, 8.0, LGRAY, t=0.04)

        tx(sl, body, cx + 0.5, cy + 5.7, 8.5, 5.8,
           size=12.5, color=SILVER)

    # Bottom insight bar
    box(sl, ML, CTT + CTH - 1.6, CW, 1.5, CARD2)
    rule(sl, ML, CTT + CTH - 1.6, CW, RED, t=0.2)
    tx(sl, "Resultado: a Capitol presta servico de alto valor, mas opera como empresa dos anos 2000.",
       ML + 0.5, CTT + CTH - 1.35, CW - 1.0, 1.1,
       size=12.5, bold=True, color=WHITE)

    return sl


def slide_04_oportunidade(prs):
    sl = interior_slide(prs, "A Oportunidade",
                        "Valor que esta sendo desperdicado hoje", page=4)

    kpis = [
        ("R$ 0",      "em receita recorrente de assinaturas digitais"),
        ("0",         "clientes captados via canais digitais no ultimo ano"),
        ("40h/mes",   "gastas em monitoramento manual que pode ser automatizado"),
        ("R$ 8-15k",  "potencial de novos contratos mensais via inbound"),
        ("3-5x",      "o alcance possivel com presenca digital estruturada"),
    ]

    row_h = CTH / len(kpis) - 0.18
    for i, (kpi, desc) in enumerate(kpis):
        cy = CTT + i * (row_h + 0.18)
        box(sl, ML, cy, CW, row_h, CARD)
        vline(sl, ML, cy, row_h, RED, t=0.4)

        tx(sl, kpi, ML + 0.75, cy + (row_h/2) - 0.75, 8.5, 1.5,
           size=26, bold=True, color=RED)
        tx(sl, desc, ML + 9.5, cy + (row_h/2) - 0.4, CW - 10.0, 1.2,
           size=13, color=SILVER)

    return sl


def slide_05_solucao(prs):
    sl = interior_slide(prs, "A Solucao LocalRise",
                        "Tres pilares integrados para transformar a operacao", page=5)

    pilars = [
        ("01", "PRESENCA\nDIGITAL",
         "Site profissional\nFunil de captacao\nConteudo automatizado\nSEO local e nacional"),
        ("02", "AUTOMACAO\nINTELIGENTE",
         "n8n como orquestrador\nAlertas em tempo real\nRelatorios automaticos\nCRM integrado"),
        ("03", "SAAS DE\nINTELIGENCIA",
         "LexRadar  |  PoliticaIQ\nRegRisk  |  OpportunityMap\nCrisisMonitor\nReceita recorrente"),
    ]

    card_w = CW / 3 - 0.3
    for i, (num, title, body) in enumerate(pilars):
        cx = ML + i * (card_w + 0.4)

        box(sl, cx, CTT, card_w, CTH, CARD)
        rule(sl, cx, CTT, card_w, RED, t=0.25)

        # Number badge
        box(sl, cx + 0.4, CTT + 0.6, 2.2, 2.0, DRED)
        tx(sl, num, cx + 0.4, CTT + 0.6, 2.2, 2.0,
           size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        tx(sl, title, cx + 0.4, CTT + 3.2, card_w - 0.8, 2.8,
           size=18, bold=True, color=WHITE)

        rule(sl, cx + 0.4, CTT + 6.1, card_w - 0.8, RED, t=0.07)

        tx(sl, body, cx + 0.4, CTT + 6.5, card_w - 0.8, CTH - 7.0,
           size=12.5, color=SILVER)

    return sl


def slide_06_automacoes(prs):
    sl = interior_slide(prs, "Automacoes Prioritarias",
                        "Processos que serao automatizados com n8n + Claude API", page=6)

    automacoes = [
        ("01", "Monitoramento Legislativo",
         "Diario  |  Alta prioridade",
         "Scraping DOU/Camara  ->  Classificacao IA  ->  Alerta WhatsApp/Email"),
        ("02", "Geracao de Relatorios",
         "Semanal/Mensal  |  Alto impacto",
         "Coleta de dados  ->  Analise Claude API  ->  Docx/PDF formatado  ->  Envio automatico"),
        ("03", "Captura de Leads",
         "Continuo  |  Estrategico",
         "Formulario site  ->  CRM  ->  Email sequence  ->  Qualificacao automatica"),
        ("04", "Agendamento de Reunioes",
         "Sob demanda  |  Operacional",
         "Calendly integrado  ->  Confirmacao automatica  ->  Briefing pre-reuniao"),
        ("05", "Gestao de Clientes",
         "Continuo  |  Recorrente",
         "Painel por cliente  ->  Alertas customizados  ->  NPS automatizado"),
    ]

    row_h = (CTH - 0.2) / len(automacoes) - 0.2
    for i, (num, title, tag, flow) in enumerate(automacoes):
        cy = CTT + i * (row_h + 0.22)
        box(sl, ML, cy, CW, row_h, CARD)
        vline(sl, ML, cy, row_h, RED, t=0.4)

        tx(sl, num, ML + 0.7, cy + (row_h/2) - 0.5, 1.5, 1.0,
           size=18, bold=True, color=DRED)
        tx(sl, title, ML + 2.3, cy + 0.18, 10, 0.9,
           size=13, bold=True, color=WHITE)
        tx(sl, tag, ML + 2.3, cy + 1.1, 10, 0.7,
           size=9.5, color=RED)
        tx(sl, flow, ML + 13.0, cy + (row_h/2) - 0.4, CW - 13.5, 1.1,
           size=11, color=SILVER)

    return sl


def slide_07_antes_depois(prs):
    sl = interior_slide(prs, "Transformacao Operacional",
                        "O antes e o depois da implantacao LocalRise", page=7)

    # Column headers
    col_w = (CW - 0.8) / 2

    box(sl, ML, CTT, col_w, 1.4, DRED)
    tx(sl, "HOJE  —  Capitol sem LocalRise",
       ML + 0.5, CTT + 0.3, col_w - 1.0, 0.9,
       size=13, bold=True, color=WHITE)

    box(sl, ML + col_w + 0.8, CTT, col_w, 1.4, RED)
    tx(sl, "APOS  —  Com LocalRise",
       ML + col_w + 1.3, CTT + 0.3, col_w - 1.0, 0.9,
       size=13, bold=True, color=WHITE)

    comparisons = [
        ("Monitoramento manual, 8-10h/semana",
         "Monitoramento automatico 24/7, alertas em minutos"),
        ("Relatorios criados do zero a cada entrega",
         "Relatorios gerados automaticamente com 1 clique"),
        ("Sem presenca digital, so indicacoes",
         "Site, funil, SEO e captacao ativa de leads"),
        ("Sem receita recorrente",
         "SaaS com assinaturas mensais previsiveis"),
        ("Clientes sem painel de acompanhamento",
         "Dashboard personalizado por cliente em tempo real"),
        ("Sem diferencial tecnologico",
         "Pioneiros em IA regulatoria no mercado brasileiro"),
    ]

    row_h = (CTH - 1.8) / len(comparisons) - 0.12
    for i, (before, after) in enumerate(comparisons):
        cy = CTT + 1.6 + i * (row_h + 0.12)

        # Before card
        box(sl, ML, cy, col_w, row_h, CARD)
        vline(sl, ML, cy, row_h, RGBColor(0x55, 0x55, 0x55), t=0.35)
        tx(sl, before, ML + 0.65, cy + (row_h/2) - 0.35, col_w - 0.9, row_h,
           size=11, color=SILVER)

        # After card
        box(sl, ML + col_w + 0.8, cy, col_w, row_h, CARD)
        vline(sl, ML + col_w + 0.8, cy, row_h, RED, t=0.35)
        tx(sl, after, ML + col_w + 1.35, cy + (row_h/2) - 0.35, col_w - 0.9, row_h,
           size=11, bold=True, color=WHITE)

    return sl


def slide_08_produtos(prs):
    sl = interior_slide(prs, "Portfolio de Produtos SaaS",
                        "Cinco linhas de receita recorrente", page=8)

    products = [
        ("LexRadar",
         "Monitoramento legislativo em tempo real. Alertas configuráveis por tema, orgao e palavra-chave.",
         "Assinatura mensal"),
        ("PoliticaIQ",
         "Analise de posicionamento de parlamentares com IA. Mapas de influencia e historico de votacoes.",
         "Assinatura + relatorio"),
        ("RegRisk",
         "Score de risco regulatorio por setor. Identifica ameacas antes que se tornem problemas.",
         "Consultivo + SaaS"),
        ("OpportunityMap",
         "Radar de licitacoes, PPPs e programas publicos alinhados ao perfil de cada cliente.",
         "Comissao + assinatura"),
        ("CrisisMonitor",
         "Monitoramento de midia e redes sociais para detectar crises regulatorias emergentes.",
         "Retainer mensal"),
    ]

    row_h = CTH / len(products) - 0.22
    for i, (name, desc, model) in enumerate(products):
        cy = CTT + i * (row_h + 0.22)
        box(sl, ML, cy, CW, row_h, CARD)
        vline(sl, ML, cy, row_h, RED, t=0.4)

        tx(sl, name, ML + 0.75, cy + 0.18, 8.5, 0.9,
           size=14, bold=True, color=RED)
        tx(sl, desc, ML + 0.75, cy + 1.1, 21, row_h - 1.3,
           size=11, color=SILVER)

        # Model badge
        box(sl, ML + CW - 7.5, cy + (row_h/2) - 0.45, 7.0, 1.1, CARD2)
        tx(sl, model, ML + CW - 7.4, cy + (row_h/2) - 0.35, 6.8, 0.9,
           size=10.5, bold=True, color=SILVER, align=PP_ALIGN.CENTER)

    return sl


def slide_09_pipeline(prs):
    sl = interior_slide(prs, "Pipeline de Inteligencia Regulatoria",
                        "Seis agentes em serie — da coleta ao cliente final", page=9)

    agents = [
        ("01", "COLLECTOR",   "Scraping DOU,\nCamara, Senado,\nTCU e portais"),
        ("02", "CLASSIFIER",  "Classifica por\ntema, urgencia\ne relevancia"),
        ("03", "ANALYST",     "Claude API\ngera analise\ntecnica"),
        ("04", "RISK SCORER", "Score 0-100\npor setor e\nperfil do cliente"),
        ("05", "REPORTER",    "Relatorio\nWord/PDF\nformatado"),
        ("06", "DISPATCHER",  "Entrega:\nEmail, WhatsApp\nou dashboard"),
    ]

    cw = (CW - 0.5 * 5) / 6   # card width
    ch = CTH - 1.0             # card height

    for i, (num, name, desc) in enumerate(agents):
        cx = ML + i * (cw + 0.5)
        cy = CTT

        box(sl, cx, cy, cw, ch, CARD)
        rule(sl, cx, cy, cw, RED, t=0.22)

        tx(sl, num, cx + 0.3, cy + 0.5, cw - 0.6, 1.4,
           size=26, bold=True, color=DRED, align=PP_ALIGN.CENTER)

        tx(sl, name, cx + 0.2, cy + 2.1, cw - 0.4, 1.6,
           size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        rule(sl, cx + 0.4, cy + 3.8, cw - 0.8, LGRAY, t=0.04)

        tx(sl, desc, cx + 0.2, cy + 4.1, cw - 0.4, ch - 4.5,
           size=10.5, color=SILVER, align=PP_ALIGN.CENTER)

        # Arrow (except last)
        if i < len(agents) - 1:
            ax = cx + cw + 0.05
            tx(sl, ">", ax, cy + (ch/2) - 0.6, 0.5, 1.2,
               size=16, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # Orchestrator tag
    box(sl, ML, CTT + ch + 0.25, CW, 0.95, CARD2)
    tx(sl, "Orquestrador: n8n  |  LLM: Claude API (Anthropic)  |  Entrega: Multi-canal",
       ML + 0.5, CTT + ch + 0.35, CW - 1.0, 0.75,
       size=11, bold=True, color=SILVER, align=PP_ALIGN.CENTER)

    return sl


def slide_10_arquitetura(prs):
    sl = interior_slide(prs, "Arquitetura Tecnologica",
                        "Stack moderno, escalavel e de baixo custo operacional", page=10)

    layers = [
        ("COLETA",        "Scrapers Python  |  RSS feeds  |  APIs publicas (Camara, Senado, DOU)  |  Web crawlers"),
        ("ORQUESTRACAO",  "n8n self-hosted  —  Workflows visuais, agendamentos, integracoes e logica de negocio"),
        ("INTELIGENCIA",  "Claude API (Anthropic)  —  Analise, classificacao, scoring e geracao de relatorios"),
        ("DADOS",         "Supabase (PostgreSQL)  —  Armazenamento, historico, perfis de clientes e logs"),
        ("ENTREGA",       "Email (SMTP)  |  WhatsApp (Evolution API)  |  Dashboard Next.js  |  PDF/Docx"),
    ]

    row_h = CTH / len(layers) - 0.2
    for i, (layer, desc) in enumerate(layers):
        cy = CTT + i * (row_h + 0.2)
        fill = CARD if i % 2 == 0 else CARD2
        box(sl, ML, cy, CW, row_h, fill)
        vline(sl, ML, cy, row_h, RED, t=0.4)

        tx(sl, layer, ML + 0.75, cy + (row_h/2) - 0.7, 7.5, 0.9,
           size=11, bold=True, color=RED)
        tx(sl, desc, ML + 8.5, cy + (row_h/2) - 0.4, CW - 9.0, 1.2,
           size=12, color=SILVER)

    return sl


def slide_11_fases(prs):
    sl = interior_slide(prs, "Plano de Implementacao",
                        "Quatro fases progressivas — do diagnostico ao produto completo", page=11)

    phases = [
        ("FASE 1", "Fundacao Digital",   "30 dias",
         "Site profissional\nSetup Google Analytics\nPerfis redes sociais\nIdentidade visual digital"),
        ("FASE 2", "Automacao Core",     "60 dias",
         "Monitoramento automatizado\nAlertas por cliente\nRelatorios automaticos\nCRM integrado"),
        ("FASE 3", "SaaS MVP",           "90 dias",
         "LexRadar em producao\nDashboard por cliente\nPrimeiros assinantes\nDocumentacao"),
        ("FASE 4", "Escala",             "180 dias",
         "Portfolio completo (5 produtos)\nTime de CS dedicado\nParcerias estrategicas\nMRR consolidado"),
    ]

    card_w = (CW - 0.4 * 3) / 4

    for i, (phase, title, deadline, items) in enumerate(phases):
        cx = ML + i * (card_w + 0.4)
        phase_card(sl, cx, CTT, card_w, CTH, phase, title, deadline, items)

        # Timeline connector
        if i < len(phases) - 1:
            rule(sl, cx + card_w, CTT + CTH/2, 0.4, RED, t=0.06)

    return sl


def slide_12_roi(prs):
    sl = interior_slide(prs, "ROI e Metricas Esperadas",
                        "Projecoes conservadoras para os primeiros 12 meses", page=12)

    kpis = [
        ("40h/mes",    "Economizadas em\nmonitoramento manual",  "liberadas para atividade estrategica"),
        ("3-5x",       "Ampliacao do alcance\nde prospeccao",    "com presenca digital estruturada"),
        ("R$ 8-15k",   "Novos contratos mensais\nvia inbound",   "captados por canal digital proprio"),
        ("R$ 5-20k",   "MRR potencial\nde assinaturas SaaS",     "no primeiro ano de operacao"),
        ("< 90 dias",  "Payback estimado\ndo investimento",       "ROI positivo no curto prazo"),
        ("10-30",      "Clientes assinantes\nno primeiro ano",    "base inicial de receita recorrente"),
    ]

    card_w = (CW - 0.4 * 2) / 3
    card_h = (CTH - 0.4) / 2

    for i, (number, label, sub) in enumerate(kpis):
        col = i % 3
        row = i // 3
        cx = ML + col * (card_w + 0.4)
        cy = CTT + row * (card_h + 0.4)
        kpi_card(sl, cx, cy, card_w, card_h, number, label, sub)

    return sl


def slide_13_modelo_comercial(prs):
    sl = interior_slide(prs, "Modelo Comercial",
                        "Tres modalidades de engajamento com a LocalRise", page=13)

    plans = [
        ("STARTER",   "Implantacao",    False,
         "Site + identidade digital\nSetup automacao basica\nPrimeiro relatorio automatico\n3 meses de suporte",
         "A partir de\nR$ 3.500"),
        ("GROWTH",    "Recomendado",    True,
         "Tudo do Starter\n+ LexRadar MVP\n+ Dashboard de clientes\n+ Suporte dedicado 6 meses",
         "A partir de\nR$ 6.800"),
        ("ENTERPRISE","Parceria Plena", False,
         "Tudo do Growth\n+ Portfolio completo\n+ White-label disponivel\n+ Revenue share SaaS",
         "Sob consulta"),
    ]

    card_w = (CW - 0.6 * 2) / 3

    for i, (name, tag, highlight, items, price) in enumerate(plans):
        cx = ML + i * (card_w + 0.6)
        cy = CTT

        fill = RED if highlight else CARD
        box(sl, cx, cy, card_w, CTH, fill)

        top_color = DRED
        rule(sl, cx, cy, card_w, top_color, t=0.4)

        if highlight:
            box(sl, cx + 1.0, cy + 0.6, card_w - 2.0, 0.9, DRED)
            tx(sl, "MAIS ESCOLHIDO", cx + 1.0, cy + 0.65, card_w - 2.0, 0.8,
               size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        ty = cy + (1.6 if highlight else 0.8)
        tx(sl, name, cx + 0.5, ty, card_w - 1.0, 1.5,
           size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tx(sl, tag, cx + 0.5, ty + 1.5, card_w - 1.0, 0.9,
           size=11, color=RGBColor(0xFF, 0xCC, 0xCC) if highlight else MUTE,
           align=PP_ALIGN.CENTER)

        rule(sl, cx + 0.7, ty + 2.5, card_w - 1.4,
             RGBColor(0xFF, 0xFF, 0xFF) if highlight else LGRAY, t=0.05)

        tx(sl, items, cx + 0.5, ty + 2.8, card_w - 1.0, CTH - ty - 3.5 + cy,
           size=12, color=WHITE, align=PP_ALIGN.LEFT)

        box(sl, cx + 0.4, cy + CTH - 2.8, card_w - 0.8, 2.4, DRED)
        tx(sl, price, cx + 0.4, cy + CTH - 2.7, card_w - 0.8, 2.2,
           size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    return sl


def slide_14_por_que_localrise(prs):
    sl = interior_slide(prs, "Por que LocalRise Advisory?",
                        "O parceiro certo para esta transformacao", page=14)

    reasons = [
        ("Especialistas em Gov-Tech",
         "Conhecimento profundo do mercado de relacoes governamentais, compliance e inteligencia regulatoria brasileira."),
        ("Stack Proprietaria",
         "Automacoes desenvolvidas internamente com n8n + Claude API. Sem dependencia de ferramentas caras de terceiros."),
        ("Entrega Orientada a Resultado",
         "Nao vendemos horas. Entregamos outcomes: leads, relatorios, alertas e receita recorrente gerados e mensurados."),
        ("Modelo de Parceria Real",
         "Alinhamento via revenue share em produtos SaaS. Nosso sucesso depende diretamente do seu sucesso."),
        ("Velocidade de Implantacao",
         "MVP rodando em 30-60 dias. Sem projetos longos de TI. Sem burocracia. Resultados visiveis rapidamente."),
    ]

    row_h = CTH / len(reasons) - 0.22
    for i, (title, body) in enumerate(reasons):
        cy = CTT + i * (row_h + 0.22)
        box(sl, ML, cy, CW, row_h, CARD)
        vline(sl, ML, cy, row_h, RED, t=0.4)

        # Number badge
        box(sl, ML + 0.65, cy + (row_h/2) - 0.55, 1.6, 1.1, DRED)
        tx(sl, f"0{i+1}", ML + 0.65, cy + (row_h/2) - 0.55, 1.6, 1.1,
           size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        tx(sl, title, ML + 2.7, cy + 0.2, 12, 0.9,
           size=13, bold=True, color=RED)
        tx(sl, body, ML + 2.7, cy + 1.1, CW - 3.2, row_h - 1.3,
           size=11.5, color=SILVER)

    return sl


def slide_15_cta(prs):
    sl = blank(prs)

    # Split background
    bg(sl, BLACK)
    box(sl, W * 0.55, 0, W * 0.45, H, RED)
    box(sl, W * 0.55 - 0.15, 0, 0.3, H, DRED)

    # ── Left: Steps ───────────────────────────────────────────────────────
    lw = W * 0.55 - 2.0

    tx(sl, "Proximos Passos", 2.0, 1.8, lw, 2.5,
       size=38, bold=True, color=WHITE)

    rule(sl, 2.0, 4.4, lw - 0.5, RED, t=0.12)

    steps = [
        "1.  Validar escopo e prioridades (esta reuniao)",
        "2.  Proposta formal com cronograma detalhado",
        "3.  Kickoff em ate 15 dias apos assinatura",
        "4.  Primeiro entregavel em ate 30 dias",
    ]
    for i, step in enumerate(steps):
        box(sl, 2.0, 5.0 + i * 2.4, lw - 0.5, 2.1, CARD)
        vline(sl, 2.0, 5.0 + i * 2.4, 2.1, RED, t=0.4)
        tx(sl, step, 2.8, 5.3 + i * 2.4, lw - 1.4, 1.5,
           size=12.5, color=SILVER)

    # Logo bottom left
    logo(sl, x=2.0, y=15.8, w=4.0)
    tx(sl, f"{DATE}",
       2.0, 15.2, 10, 0.8,
       size=9, color=MUTE)

    # ── Right: Contact ────────────────────────────────────────────────────
    rx = W * 0.55 + 0.8
    rw = W * 0.45 - 1.6

    tx(sl, "Vamos\nconstruir\njuntos?",
       rx, 2.0, rw, 6.5,
       size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    rule(sl, rx, 8.5, rw - 0.5, RGBColor(0xFF, 0xFF, 0xFF), t=0.08)

    tx(sl, "LocalRise Advisory", rx, 9.0, rw, 1.2,
       size=16, bold=True, color=WHITE)
    tx(sl, "Nicolas Slater", rx, 10.2, rw, 0.9,
       size=13, color=WHITE)
    tx(sl, "Diretoria Comercial", rx, 11.0, rw, 0.8,
       size=11, color=RGBColor(0xFF, 0xCC, 0xCC))

    rule(sl, rx, 12.2, rw - 0.5, RGBColor(0xFF, 0xFF, 0xFF), t=0.05)

    tx(sl, "contato@localrise.com.br", rx, 12.6, rw, 0.9,
       size=12, color=WHITE)
    tx(sl, "localrise.com.br", rx, 13.4, rw, 0.9,
       size=12, bold=True, color=WHITE)

    return sl


# ════════════════════════════════════════════════════════════════════════════
# BUILD
# ════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    builders = [
        (slide_01_cover,              "Capa"),
        (slide_02_sumario,            "Sumario Executivo"),
        (slide_03_diagnostico,        "Diagnostico"),
        (slide_04_oportunidade,       "Oportunidade"),
        (slide_05_solucao,            "Solucao"),
        (slide_06_automacoes,         "Automacoes"),
        (slide_07_antes_depois,       "Antes x Depois"),
        (slide_08_produtos,           "Produtos SaaS"),
        (slide_09_pipeline,           "Pipeline de IA"),
        (slide_10_arquitetura,        "Arquitetura"),
        (slide_11_fases,              "Fases"),
        (slide_12_roi,                "ROI"),
        (slide_13_modelo_comercial,   "Modelo Comercial"),
        (slide_14_por_que_localrise,  "Por que LocalRise"),
        (slide_15_cta,                "CTA Final"),
    ]

    print("Gerando apresentacao premium v2...")
    for fn, name in builders:
        fn(prs)
        print(f"  OK  {name}")

    prs.save(str(OUT_PATH))
    print("\nPronto!")
    print("  " + str(OUT_PATH))


if __name__ == "__main__":
    build()
