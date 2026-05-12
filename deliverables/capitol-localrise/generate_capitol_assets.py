from pathlib import Path

from docx import Document
from docx.enum.table import WD_TABLE_ALIGNMENT, WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_LINE_SPACING
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Cm, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PPTColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm as PCm, Pt as PPt


BASE_DIR = Path(__file__).resolve().parent
LOGO_PATH = Path(
    r"C:\Users\digui\Documents\localrise-brain\Claude\Claude Local RIse\assents\logo Local Rise Principal.png"
)

RED = RGBColor(0xE3, 0x1B, 0x23)
BLACK = RGBColor(0x0A, 0x0A, 0x0A)
GRAY = RGBColor(0x5F, 0x63, 0x68)
LIGHT = RGBColor(0xF6, 0xF7, 0xF8)

PPT_RED = PPTColor(0xE3, 0x1B, 0x23)
PPT_BLACK = PPTColor(0x0A, 0x0A, 0x0A)
PPT_WHITE = PPTColor(0xFF, 0xFF, 0xFF)
PPT_SOFT = PPTColor(0xE8, 0xEA, 0xED)

COMPANY = "Capitol Consultoria"
HANDLE = "@capitolconsultoria"
DATE = "08 de abril de 2026"


def set_cell_border(cell, **kwargs):
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()
    tc_borders = tc_pr.first_child_found_in("w:tcBorders")
    if tc_borders is None:
        tc_borders = OxmlElement("w:tcBorders")
        tc_pr.append(tc_borders)
    for edge in ("left", "top", "right", "bottom"):
        edge_data = kwargs.get(edge)
        if not edge_data:
            continue
        tag = "w:%s" % edge
        element = tc_borders.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            tc_borders.append(element)
        for key in ("val", "sz", "space", "color"):
            if key in edge_data:
                element.set(qn("w:%s" % key), str(edge_data[key]))


def apply_doc_style(document: Document):
    section = document.sections[0]
    section.page_width = Cm(21)
    section.page_height = Cm(29.7)
    section.top_margin = Cm(2.2)
    section.bottom_margin = Cm(2.2)
    section.left_margin = Cm(2.2)
    section.right_margin = Cm(2.2)

    normal = document.styles["Normal"]
    normal.font.name = "Calibri"
    normal._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
    normal.font.size = Pt(11)
    normal.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    normal.paragraph_format.space_after = Pt(8)


def add_brand_header(document: Document, title: str, subtitle: str | None = None):
    if LOGO_PATH.exists():
        paragraph = document.add_paragraph()
        paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = paragraph.add_run()
        run.add_picture(str(LOGO_PATH), width=Cm(5.4))

    brand = document.add_paragraph()
    brand.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = brand.add_run("LOCALRISE ADVISORY")
    run.bold = True
    run.font.name = "Calibri"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
    run.font.size = Pt(13)
    run.font.color.rgb = RED

    table = document.add_table(rows=1, cols=1)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    cell = table.cell(0, 0)
    set_cell_border(cell, bottom={"val": "single", "sz": 8, "space": 0, "color": "D9DDE3"})
    cell.text = ""

    title_p = document.add_paragraph()
    title_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_p.paragraph_format.space_before = Pt(10)
    title_run = title_p.add_run(title)
    title_run.bold = True
    title_run.font.name = "Calibri"
    title_run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
    title_run.font.size = Pt(18)

    if subtitle:
        subtitle_p = document.add_paragraph()
        subtitle_p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        subtitle_run = subtitle_p.add_run(subtitle)
        subtitle_run.font.name = "Calibri"
        subtitle_run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        subtitle_run.font.size = Pt(10.5)
        subtitle_run.font.color.rgb = GRAY


def add_section_title(document: Document, text: str):
    p = document.add_paragraph()
    p.paragraph_format.space_before = Pt(10)
    p.paragraph_format.space_after = Pt(4)
    run = p.add_run(text.upper())
    run.bold = True
    run.font.name = "Calibri"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
    run.font.size = Pt(12)
    run.font.color.rgb = RED


def add_paragraph(document: Document, text: str, bold_prefix: str | None = None):
    p = document.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
    p.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
    if bold_prefix and text.startswith(bold_prefix):
        run = p.add_run(bold_prefix)
        run.bold = True
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(11)
        run2 = p.add_run(text[len(bold_prefix):])
        run2.font.name = "Calibri"
        run2._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run2.font.size = Pt(11)
    else:
        run = p.add_run(text)
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(11)


def add_bullets(document: Document, items: list[str]):
    for item in items:
        p = document.add_paragraph(style="List Bullet")
        p.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
        run = p.add_run(item)
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(11)


def add_numbered(document: Document, items: list[str]):
    for item in items:
        p = document.add_paragraph(style="List Number")
        p.paragraph_format.line_spacing_rule = WD_LINE_SPACING.ONE_POINT_FIVE
        run = p.add_run(item)
        run.font.name = "Calibri"
        run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
        run.font.size = Pt(11)


def add_signature_block(document: Document):
    p = document.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = p.add_run(f"Sao Paulo, {DATE}.")
    run.bold = True
    run.font.name = "Calibri"
    run._element.rPr.rFonts.set(qn("w:eastAsia"), "Calibri")
    run.font.size = Pt(11)

    table = document.add_table(rows=1, cols=2)
    table.alignment = WD_TABLE_ALIGNMENT.CENTER
    for row in table.rows:
        for cell in row.cells:
            cell.width = Cm(7.5)
            cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.BOTTOM
            set_cell_border(cell, top={"val": "single", "sz": 8, "space": 0, "color": "AEB4BB"})
    left = table.cell(0, 0).paragraphs[0]
    left.alignment = WD_ALIGN_PARAGRAPH.CENTER
    left.add_run("LOCALRISE ADVISORY\n").bold = True
    left.add_run("Nicolas Slater\n")
    left.add_run("Diretoria Comercial")
    right = table.cell(0, 1).paragraphs[0]
    right.alignment = WD_ALIGN_PARAGRAPH.CENTER
    right.add_run(f"{COMPANY}\n").bold = True
    right.add_run("Responsavel: ______________________\n")
    right.add_run("Cargo: ___________________________")


def save_proposal_doc():
    doc = Document()
    apply_doc_style(doc)
    add_brand_header(
        doc,
        "Proposta Comercial Estrategica",
        "Capitol Consultoria | site, Instagram e Meta Ads",
    )

    add_section_title(doc, "Contexto")
    add_paragraph(
        doc,
        "Esta proposta parte do diagnostico estrategico realizado pela LocalRise Advisory em 08/04/2026. A leitura central e que a Capitol possui senioridade real e tema de alta relevancia, mas ainda nao converte essa autoridade em uma presenca digital capaz de gerar confianca, recorrencia e demanda previsivel.",
    )
    add_paragraph(
        doc,
        "Fatos observados: perfil publico no Instagram com baixa tracao relativa (87 seguidores, 28 posts) e ausencia de site oficial claro/indexado na verificacao publica. Hipotese estrategica: a marca esta subposicionada digitalmente em relacao ao nivel de senioridade que entrega no offline.",
    )

    add_section_title(doc, "O que encontramos")
    add_bullets(
        doc,
        [
            "Instagram existente, mas hoje mais proximo de cartao de visita do que de canal de autoridade e aquisicao.",
            "Ausencia de site funcional/publicamente validado, o que enfraquece autoridade, captura de lead, mensuracao e remarketing.",
            "Posicionamento correto no tema, mas ainda pouco claro quanto ao diferencial, prova e chamada comercial.",
            "Meta Ads faz sentido apenas depois da estruturacao da base digital de conversao.",
        ],
    )

    add_section_title(doc, "Objetivo da entrada LocalRise")
    add_paragraph(
        doc,
        "Estruturar uma presenca digital coerente com o nivel institucional da Capitol, criando uma base capaz de sustentar autoridade, gerar leads qualificados e preparar a marca para escalar aquisicao com menor desperdicio de verba.",
    )

    add_section_title(doc, "Escopo recomendado")
    add_numbered(
        doc,
        [
            "Fase 1 | Site institucional-comercial com arquitetura de conversao, provas, formulacao de mensagem e preparacao para midia.",
            "Fase 2 | Reposicionamento do Instagram como canal de autoridade, aquecimento e validacao social.",
            "Fase 3 | Estrutura de Meta Ads com foco em awareness qualificado, remarketing e captacao de leads.",
        ],
    )

    add_section_title(doc, "Entregas principais")
    add_bullets(
        doc,
        [
            "Definicao de mensagem central, proposta de valor e narrativa comercial para a Capitol.",
            "Site com paginas estrategicas, formularios, CTA de reuniao e integracoes de mensuracao.",
            "Reestruturacao de bio, destaques, linha editorial e rotina de conteudo para Instagram.",
            "Planejamento de campanha e funil para Meta Ads, ativado apenas quando a base estiver pronta.",
        ],
    )

    add_section_title(doc, "Plano 30 60 90 dias")
    add_bullets(
        doc,
        [
            "0-30 dias: reposicionamento, definicao de ICP, copy central, sitemap, wireframe e novo setup do Instagram.",
            "31-60 dias: publicacao do site, ativacao de CTA, conteudo de autoridade e ajustes finos de conversao.",
            "61-90 dias: implementacao de pixel/CAPI, landing pages, remarketing e inicio da camada de midia.",
        ],
    )

    add_section_title(doc, "Resultados esperados")
    add_bullets(
        doc,
        [
            "Melhora imediata de percepcao institucional e confianca comercial.",
            "Aumento da capacidade de transformar interesse em contato qualificado.",
            "Base pronta para receber trafego pago com maior chance de retorno.",
            "Maior clareza para conversas comerciais, propostas e validacao de marca.",
        ],
    )

    add_section_title(doc, "Modelo recomendado")
    add_paragraph(doc, "Investimento da Fase 1: ________________________________", bold_prefix="Investimento da Fase 1:")
    add_paragraph(doc, "Investimento mensal das Fases 2 e 3: ________________________________", bold_prefix="Investimento mensal das Fases 2 e 3:")
    add_paragraph(
        doc,
        "Observacao: os campos de investimento foram deixados em aberto para definicao comercial interna da LocalRise antes do envio final ao cliente.",
    )

    add_section_title(doc, "Proximos passos")
    add_numbered(
        doc,
        [
            "Validar internamente a proposta de valor e o escopo inicial.",
            "Apresentar o diagnostico e a ordem recomendada de entrada em reuniao.",
            "Fechar o escopo da Fase 1 e alinhar cronograma de kickoff.",
        ],
    )

    add_signature_block(doc)
    doc.save(BASE_DIR / "proposta-comercial-capitol-localrise.docx")


def save_instagram_plan_doc():
    doc = Document()
    apply_doc_style(doc)
    add_brand_header(
        doc,
        "Plano Inicial de Instagram",
        "Capitol Consultoria | linha editorial, bio e rotina de autoridade",
    )

    add_section_title(doc, "Leitura de partida")
    add_paragraph(
        doc,
        "Fatos observados: perfil publico ativo, 87 seguidores, 54 seguindo, 28 posts e bio resumida como consultoria especializada em relacoes institucionais e governamentais. Hipotese estrategica: o perfil ainda nao traduz senioridade em desejo, prova e conversao.",
    )

    add_section_title(doc, "Objetivo do canal")
    add_bullets(
        doc,
        [
            "Transformar o Instagram em vitrine de autoridade, nao apenas em vitrine institucional.",
            "Aquecer mercado e reduzir friccao antes de contato comercial.",
            "Criar percepcao de consistencia, senioridade e capacidade analitica.",
        ],
    )

    add_section_title(doc, "Bio recomendada")
    add_bullets(
        doc,
        [
            "Linha 1: Relacoes governamentais e institucionais para empresas que precisam navegar Brasilia com clareza.",
            "Linha 2: Analise, articulacao e leitura estrategica para setores regulados e agendas sensiveis.",
            "Linha 3: CTA objetivo: agende uma conversa estrategica.",
            "Link: landing page ou formulario de contato, nao perfil solto.",
        ],
    )

    add_section_title(doc, "Destaques recomendados")
    add_bullets(
        doc,
        [
            "Sobre",
            "Setores",
            "Servicos",
            "Insights",
            "Equipe",
            "Contato",
        ],
    )

    add_section_title(doc, "Pilares de conteudo")
    add_numbered(
        doc,
        [
            "Autoridade de mercado: leitura de cenario, movimentos de Brasilia, impactos para empresas.",
            "Educativo executivo: explicar temas complexos em linguagem de decisor.",
            "Prova e metodo: como a Capitol pensa, estrutura e acompanha agendas sensiveis.",
            "Marca e bastidores: senioridade da equipe, postura institucional, agenda e bastidores sem perder sobriedade.",
            "Conversao: convite para diagnostico, conversa estrategica e contatos qualificados.",
        ],
    )

    add_section_title(doc, "Frequencia ideal")
    add_bullets(
        doc,
        [
            "3 posts por semana como base minima.",
            "2 reels curtos por semana com leitura executiva ou comentario de cenario.",
            "Stories em dias uteis com bastidores, agenda, cortes e reforco de autoridade.",
        ],
    )

    add_section_title(doc, "Calendario inicial de 30 dias")
    add_bullets(
        doc,
        [
            "Semana 1: manifesto de posicionamento, carrossel sobre o papel de relacoes governamentais, reel com leitura de cenario, story com bastidor.",
            "Semana 2: post sobre risco de operar sem leitura institucional, carrossel setorial, reel sobre erro comum de empresas em Brasilia.",
            "Semana 3: conteudo de prova/metodo, insight regulatorio da semana, reel com comentario direto do socio.",
            "Semana 4: caso ilustrativo sem expor cliente, post de perguntas frequentes, CTA para conversa estrategica.",
        ],
    )

    add_section_title(doc, "Tipos de criativo")
    add_bullets(
        doc,
        [
            "Carrosseis com tese + implicacao de negocio.",
            "Reels com camera direta, tom sobrio e hook executivo.",
            "Posts institucionais com fotografia limpa e titulacao forte.",
            "Stories com blocos curtos: contexto, implicacao, convite.",
        ],
    )

    add_section_title(doc, "Metricas prioritarias")
    add_bullets(
        doc,
        [
            "Cliques no link da bio.",
            "Salvamentos e compartilhamentos por post.",
            "Respostas em story e DMs qualificadas.",
            "Visitantes que avancam para o site ou formulario.",
        ],
    )

    add_signature_block(doc)
    doc.save(BASE_DIR / "plano-instagram-capitol.docx")


def save_site_structure_doc():
    doc = Document()
    apply_doc_style(doc)
    add_brand_header(
        doc,
        "Estrutura Ideal do Site",
        "Capitol Consultoria | sitemap, CRO e copy-base",
    )

    add_section_title(doc, "Objetivo do site")
    add_paragraph(
        doc,
        "Criar um ativo digital que sustente a senioridade institucional da Capitol, funcione como prova de autoridade e permita captacao clara de oportunidades comerciais.",
    )

    add_section_title(doc, "Arquitetura recomendada")
    add_numbered(
        doc,
        [
            "Home",
            "Quem Somos",
            "Servicos",
            "Setores e Temas",
            "Insights",
            "Contato",
        ],
    )

    add_section_title(doc, "Home | estrutura")
    add_bullets(
        doc,
        [
            "Hero com headline clara, subheadline e CTA para conversa estrategica.",
            "Bloco de credibilidade com experiencia, setores atendidos e enfoque em Brasilia.",
            "Bloco de servicos com 3 a 4 frentes principais.",
            "Bloco de metodologia.",
            "Bloco de insights/prova.",
            "CTA final com formulario curto.",
        ],
    )

    add_section_title(doc, "Copy-base da Home")
    add_paragraph(
        doc,
        "Headline sugerida: Clareza institucional para empresas que precisam decidir melhor em Brasilia.",
    )
    add_paragraph(
        doc,
        "Subheadline sugerida: A Capitol conecta leitura de cenario, articulacao institucional e acompanhamento estrategico para empresas que operam em ambientes regulatorios e politicos complexos.",
    )
    add_paragraph(
        doc,
        "CTA principal sugerido: Agendar conversa estrategica.",
    )

    add_section_title(doc, "Pagina Servicos")
    add_bullets(
        doc,
        [
            "Assessoria em relacoes governamentais e institucionais.",
            "Monitoramento e leitura de cenario.",
            "Acompanhamento de agendas e stakeholders.",
            "Subsidiacao tecnica e posicionamento institucional.",
        ],
    )

    add_section_title(doc, "Pagina Setores e Temas")
    add_bullets(
        doc,
        [
            "Setores regulados e agendas sensiveis.",
            "Relacao entre tema, risco e impacto de negocio.",
            "Conteudo pensado para decisores, nao para publico generico.",
        ],
    )

    add_section_title(doc, "Elementos obrigatorios de conversao")
    add_bullets(
        doc,
        [
            "CTA acima da dobra e repetido ao longo da pagina.",
            "Formulario simples com nome, empresa, cargo, e-mail e mensagem.",
            "Botao para WhatsApp apenas se o atendimento comercial for rapido.",
            "Bloco de prova social ou credenciais institucionais.",
            "Pixel Meta, GA4, Search Console e mapas de conversao desde o primeiro dia.",
        ],
    )

    add_section_title(doc, "SEO basico necessario")
    add_bullets(
        doc,
        [
            "Titles e metas focados em marca + relacoes governamentais + Brasilia.",
            "Estrutura H1/H2 clara por pagina.",
            "URLs limpas.",
            "Schema Organization e Article para insights.",
            "Conteudo indexavel para marca e temas estrategicos.",
        ],
    )

    add_section_title(doc, "Preparo para Meta Ads")
    add_bullets(
        doc,
        [
            "Landing page dedicada para campanhas futuras.",
            "Mensagem coerente entre anuncio e pagina.",
            "Oferta clara: conversa estrategica, diagnostico ou material executivo.",
            "Eventos de conversao configurados antes de escalar verba.",
        ],
    )

    add_signature_block(doc)
    doc.save(BASE_DIR / "estrutura-site-capitol.docx")


def add_slide_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(PCm(0.9), PCm(0.8), PCm(11.5), PCm(1.3))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Calibri"
    run.font.bold = True
    run.font.size = PPt(24)
    run.font.color.rgb = PPT_WHITE
    if subtitle:
        sub_box = slide.shapes.add_textbox(PCm(0.95), PCm(2.0), PCm(11.2), PCm(0.8))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.name = "Calibri"
        p2.font.size = PPt(11)
        p2.font.color.rgb = PPT_SOFT


def add_slide_bullets(slide, items, left=0.95, top=2.9, width=11.2, height=4.0):
    box = slide.shapes.add_textbox(PCm(left), PCm(top), PCm(width), PCm(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Calibri"
        p.font.size = PPt(17)
        p.font.color.rgb = PPT_WHITE
        p.bullet = True
        p.space_after = PPt(8)


def style_slide(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PPT_BLACK
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PCm(0), PCm(0), PCm(13.33), PCm(0.45))
    band.fill.solid()
    band.fill.fore_color.rgb = PPT_RED
    band.line.fill.background()
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), PCm(10.8), PCm(0.65), width=PCm(1.8))


def save_presentation():
    prs = Presentation()
    prs.slide_width = PCm(33.867)
    prs.slide_height = PCm(19.05)

    slides = []
    for _ in range(10):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        style_slide(slide)
        slides.append(slide)

    add_slide_title(slides[0], "Capitol Consultoria", "Diagnostico estrategico e proposta de entrada LocalRise | abril 2026")
    cover_box = slides[0].shapes.add_textbox(PCm(0.95), PCm(4.0), PCm(8.8), PCm(2.5))
    p = cover_box.text_frame.paragraphs[0]
    p.text = "Como transformar senioridade institucional em presenca digital, autoridade escalavel e aquisicao preparada para crescer."
    p.font.name = "Calibri"
    p.font.size = PPt(26)
    p.font.bold = True
    p.font.color.rgb = PPT_WHITE

    add_slide_title(slides[1], "O que encontramos")
    add_slide_bullets(
        slides[1],
        [
            "Instagram existe, mas ainda nao opera como ativo de autoridade e aquisicao.",
            "Nao encontramos site oficial claro e funcional na verificacao publica.",
            "A marca parece mais forte no offline do que no digital.",
            "Meta Ads agora teria risco alto de desperdiccio de verba.",
        ],
    )

    add_slide_title(slides[2], "Riscos de continuar como esta")
    add_slide_bullets(
        slides[2],
        [
            "Baixa validacao institucional para novos decisores.",
            "Pouca capacidade de captar, medir e qualificar leads.",
            "Dependencia excessiva de reputacao offline e networking.",
            "Trafego pago sem base de conversao tende a performar mal.",
        ],
    )

    add_slide_title(slides[3], "A ordem ideal de entrada")
    add_slide_bullets(
        slides[3],
        [
            "1. Site institucional-comercial",
            "2. Instagram de autoridade",
            "3. Meta Ads como acelerador, nao como ponto de partida",
        ],
    )

    add_slide_title(slides[4], "Frente 1 | Site")
    add_slide_bullets(
        slides[4],
        [
            "Criar base de autoridade, prova e conversao.",
            "Arquitetura recomendada: Home, Quem Somos, Servicos, Setores, Insights e Contato.",
            "CTA para conversa estrategica e mensuracao pronta para midia.",
        ],
    )

    add_slide_title(slides[5], "Frente 2 | Instagram")
    add_slide_bullets(
        slides[5],
        [
            "Reposicionar o perfil para transmitir senioridade, clareza e consistencia.",
            "Pilares: autoridade, cenario, metodo, bastidores e conversao.",
            "Calendario minimo: 3 posts por semana + stories uteis em dias uteis.",
        ],
    )

    add_slide_title(slides[6], "Frente 3 | Meta Ads")
    add_slide_bullets(
        slides[6],
        [
            "Entrar depois do site e da oferta estarem prontos.",
            "Logica inicial: awareness qualificado, conteudo, remarketing e captacao.",
            "Sem base pronta, a campanha tende a comprar atencao e perder conversao.",
        ],
    )

    add_slide_title(slides[7], "Plano 90 dias")
    add_slide_bullets(
        slides[7],
        [
            "30 dias: posicionamento, copy central, sitemap, bio e ativos basicos.",
            "60 dias: site publicado, CTA funcionando e linha editorial rodando.",
            "90 dias: pixel, landing pages e ativacao de midia com criterio.",
        ],
    )

    add_slide_title(slides[8], "O que a LocalRise entrega")
    add_slide_bullets(
        slides[8],
        [
            "Mensagem comercial mais clara.",
            "Ativo proprio de captacao e prova.",
            "Conteudo de autoridade que prepara reunioes.",
            "Estrutura para escalar aquisicao com menos desperdicio.",
        ],
    )

    add_slide_title(slides[9], "Proximo passo recomendado")
    add_slide_bullets(
        slides[9],
        [
            "Validar internamente a Fase 1 como primeira venda.",
            "Levar esta narrativa para a reuniao com foco em site primeiro.",
            "Deixar Instagram e Meta Ads como expansao natural da base criada.",
        ],
    )

    prs.save(BASE_DIR / "apresentacao-executiva-capitol.pptx")


def main():
    BASE_DIR.mkdir(parents=True, exist_ok=True)
    save_proposal_doc()
    save_instagram_plan_doc()
    save_site_structure_doc()
    save_presentation()


if __name__ == "__main__":
    main()
