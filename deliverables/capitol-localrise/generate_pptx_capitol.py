"""
generate_pptx_capitol.py
Gera apresentação executiva Capitol Consultoria × LocalRise Advisory
em formato PowerPoint (.pptx) com design premium, 16:9, 15 slides.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Emu, Pt

# ─── Paths ─────────────────────────────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parent
LOGO_PATH = Path(
    r"C:\Users\digui\Documents\localrise-brain\Claude\Claude Local RIse\assents\logo Local Rise Principal.png"
)
OUT_PATH = BASE_DIR / "apresentacao-capitol-localrise.pptx"

# ─── Brand Palette ──────────────────────────────────────────────────────────
RED    = RGBColor(0xE3, 0x1B, 0x23)
BLACK  = RGBColor(0x0A, 0x0A, 0x0A)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x5F, 0x63, 0x68)
LGRAY  = RGBColor(0xE8, 0xEA, 0xED)
DRED   = RGBColor(0xB5, 0x12, 0x18)   # dark red accent
MBLACK = RGBColor(0x1C, 0x1C, 0x1C)   # medium black for cards

# ─── Slide dimensions (16:9 fullHD) ─────────────────────────────────────────
W = Cm(33.867)
H = Cm(19.05)

# ─── Constants ───────────────────────────────────────────────────────────────
DATE    = "08 de abril de 2026"
COMPANY = "Capitol Consultoria"


# ═══════════════════════════════════════════════════════════════════════════
# Low-level helpers
# ═══════════════════════════════════════════════════════════════════════════

def _rgb(r, g, b):
    return RGBColor(r, g, b)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    """Add a truly blank slide (no placeholders)."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


# ─── Shape helpers ───────────────────────────────────────────────────────────

def rect(slide, x, y, w, h, fill: RGBColor, alpha=None):
    """Solid filled rectangle."""
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    shape.line.width = 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape


def line_sep(slide, x, y, w, color: RGBColor, thickness=0.06):
    """Thin horizontal rule."""
    s = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(thickness))
    s.line.fill.background()
    s.line.width = 0
    s.fill.solid()
    s.fill.fore_color.rgb = color
    return s


def add_logo(slide, x=0.6, y=0.4, w=3.4):
    if LOGO_PATH.exists():
        slide.shapes.add_picture(str(LOGO_PATH), Cm(x), Cm(y), width=Cm(w))


# ─── Text helpers ────────────────────────────────────────────────────────────

def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, font="Calibri",
          wrap=True, valign=None):
    """Add a transparent text box."""
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    if valign:
        from pptx.enum.text import MSO_ANCHOR
        tf.auto_size = None
        tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def txbox_multi(slide, lines, x, y, w, h,
                size=16, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, font="Calibri",
                line_spacing=None):
    """Add text box with multiple paragraphs."""
    box = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        if line_spacing:
            from pptx.util import Pt as PPt
            from pptx.oxml.ns import qn
            from lxml import etree
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            lnSpc = etree.SubElement(pPr, qn('a:lnSpc'))
            spcPts = etree.SubElement(lnSpc, qn('a:spcPts'))
            spcPts.set('val', str(int(line_spacing * 100)))
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


# ═══════════════════════════════════════════════════════════════════════════
# Reusable slide components
# ═══════════════════════════════════════════════════════════════════════════

def section_header_bar(slide, title, subtitle=None):
    """Black left column + red accent strip for interior slides."""
    # Full black background
    rect(slide, 0, 0, 33.867, 19.05, BLACK)
    # Red vertical stripe on left
    rect(slide, 0, 0, 0.55, 19.05, RED)
    # Top thin red line
    rect(slide, 0.55, 0, 33.3, 0.22, DRED)
    # Logo top-right
    add_logo(slide, x=29.8, y=0.3, w=3.4)
    # Title
    txbox(slide, title,
          x=1.4, y=0.6, w=26, h=2.2,
          size=32, bold=True, color=WHITE)
    if subtitle:
        txbox(slide, subtitle,
              x=1.4, y=2.8, w=26, h=1.1,
              size=15, bold=False, color=LGRAY)


def footer_bar(slide, note="Confidencial | LocalRise Advisory × Capitol Consultoria | Abril 2026"):
    rect(slide, 0, 17.85, 33.867, 1.2, MBLACK)
    txbox(slide, note,
          x=1.0, y=18.05, w=22, h=0.9,
          size=9, color=GRAY, align=PP_ALIGN.LEFT)
    txbox(slide, "localrise.com.br",
          x=25, y=18.05, w=8, h=0.9,
          size=9, color=RED, align=PP_ALIGN.RIGHT)


def card(slide, x, y, w, h, fill=MBLACK, border_color=None):
    """Rounded card block."""
    s = rect(slide, x, y, w, h, fill)
    if border_color:
        s.line.color.rgb = border_color
        from pptx.util import Pt as PPt
        s.line.width = PPt(1)
    return s


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def slide_01_cover(prs):
    """Slide 1 – Capa impactante"""
    sl = blank_slide(prs)

    # Full black background
    rect(sl, 0, 0, 33.867, 19.05, BLACK)

    # Big red diagonal accent block (right side)
    rect(sl, 22.5, 0, 11.367, 19.05, RED)

    # Dark overlay band center
    rect(sl, 20.5, 0, 3.5, 19.05, DRED)

    # Left content area
    txbox(sl, "INTELIGÊNCIA",
          x=1.5, y=1.8, w=20, h=3.5,
          size=58, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txbox(sl, "REGULATÓRIA",
          x=1.5, y=4.8, w=20, h=3.5,
          size=58, bold=True, color=RED, align=PP_ALIGN.LEFT)
    txbox(sl, "COMO VANTAGEM COMPETITIVA",
          x=1.5, y=7.9, w=20, h=1.6,
          size=20, bold=False, color=LGRAY, align=PP_ALIGN.LEFT)

    # Separator line
    line_sep(sl, 1.5, 9.8, 18, WHITE, thickness=0.1)

    # Sub-info
    txbox(sl, f"Proposta Estratégica para {COMPANY}",
          x=1.5, y=10.2, w=20, h=1.1,
          size=14, bold=False, color=LGRAY)
    txbox(sl, f"Apresentada por LocalRise Advisory  ·  {DATE}",
          x=1.5, y=11.3, w=20, h=1.0,
          size=12, bold=False, color=GRAY)

    # Right column labels
    txbox(sl, "AUTOMAÇÃO\nINTELIGENTE",
          x=23.2, y=3.0, w=9, h=3,
          size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, "MONITORAMENTO\nEM TEMPO REAL",
          x=23.2, y=7.0, w=9, h=3,
          size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, "PRODUTOS\nSAAS B2B",
          x=23.2, y=11.0, w=9, h=3,
          size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Logo bottom-left
    add_logo(sl, x=1.5, y=16.2, w=4.0)

    return sl


def slide_02_sumario(prs):
    """Slide 2 – Sumário Executivo"""
    sl = blank_slide(prs)
    section_header_bar(sl, "Sumário Executivo",
                       "O que vamos apresentar hoje")
    footer_bar(sl)

    topics = [
        ("01", "Diagnóstico",      "Onde a Capitol está hoje"),
        ("02", "Oportunidade",     "O que está sendo deixado na mesa"),
        ("03", "A Solução",        "Automação + IA + SaaS"),
        ("04", "Produtos",         "LexRadar, PoliticaIQ e mais"),
        ("05", "Pipeline de IA",   "6 agentes, 1 sistema"),
        ("06", "Implementação",    "Fases e cronograma"),
        ("07", "ROI",              "Retorno esperado"),
        ("08", "Proposta Comercial","Investimento e próximos passos"),
    ]

    col_w = 14.0
    row_h = 2.0
    for i, (num, title, desc) in enumerate(topics):
        col = i % 2
        row = i // 2
        cx = 1.4 + col * 16.0
        cy = 4.0 + row * 2.6

        card(sl, cx, cy, col_w, row_h + 0.4, fill=MBLACK)
        txbox(sl, num,
              x=cx + 0.3, y=cy + 0.2, w=1.5, h=row_h,
              size=28, bold=True, color=RED)
        txbox(sl, title,
              x=cx + 2.0, y=cy + 0.15, w=11, h=0.9,
              size=15, bold=True, color=WHITE)
        txbox(sl, desc,
              x=cx + 2.0, y=cy + 1.0, w=11, h=0.8,
              size=11, color=GRAY)

    return sl


def slide_03_diagnostico(prs):
    """Slide 3 – Diagnóstico: situação atual"""
    sl = blank_slide(prs)
    section_header_bar(sl, "Diagnóstico",
                       "Onde a Capitol Consultoria está hoje")
    footer_bar(sl)

    # 3 problem cards across
    problems = [
        ("Presença Digital Mínima",
         "87 seguidores no Instagram.\nSem site. Sem blog.\nSem funil de captação."),
        ("Dependência de Indicações",
         "100% dos clientes via networking.\nNenhum canal ativo de prospecção.\nRisco alto de sazonalidade."),
        ("Operação Manual",
         "Monitoramento legislativo feito à mão.\nAlertas não automatizados.\nHoras perdidas em tarefas repetitivas."),
    ]

    for i, (title, body) in enumerate(problems):
        cx = 1.4 + i * 10.8
        cy = 4.2
        card(sl, cx, cy, 10.0, 8.0, fill=MBLACK, border_color=DRED)
        # Red top accent
        rect(sl, cx, cy, 10.0, 0.35, RED)
        txbox(sl, title,
              x=cx + 0.4, y=cy + 0.7, w=9.2, h=1.6,
              size=16, bold=True, color=WHITE)
        txbox(sl, body,
              x=cx + 0.4, y=cy + 2.4, w=9.2, h=5.2,
              size=13, color=LGRAY)

    # Bottom insight
    rect(sl, 1.4, 13.0, 31.0, 1.6, DRED)
    txbox(sl, "⚡  Resultado: a Capitol presta consultoria de alto valor, mas opera como empresa dos anos 2000.",
          x=2.0, y=13.15, w=30, h=1.3,
          size=14, bold=True, color=WHITE)

    return sl


def slide_04_oportunidade(prs):
    """Slide 4 – O que está sendo deixado na mesa"""
    sl = blank_slide(prs)
    section_header_bar(sl, "A Oportunidade",
                       "O que está sendo deixado na mesa")
    footer_bar(sl)

    items = [
        ("R$ 0", "em receita recorrente de assinaturas digitais"),
        ("0", "clientes captados via canais digitais no último ano"),
        ("40h/mês", "gastas em monitoramento manual que pode ser automatizado"),
        ("3–5×", "o alcance possível com uma presença digital estruturada"),
        ("R$ 8–15k/mês", "potencial de novos contratos via inbound marketing"),
    ]

    for i, (kpi, desc) in enumerate(items):
        cy = 4.0 + i * 2.5
        card(sl, 1.4, cy, 31.0, 2.1, fill=MBLACK)
        rect(sl, 1.4, cy, 0.45, 2.1, RED)
        txbox(sl, kpi,
              x=2.3, y=cy + 0.3, w=7, h=1.5,
              size=24, bold=True, color=RED)
        txbox(sl, desc,
              x=9.5, y=cy + 0.45, w=23, h=1.2,
              size=15, color=LGRAY)

    return sl


def slide_05_solucao(prs):
    """Slide 5 – A Solução: três pilares"""
    sl = blank_slide(prs)
    section_header_bar(sl, "A Solução LocalRise",
                       "Três pilares integrados para transformar a operação")
    footer_bar(sl)

    pilars = [
        ("01", "PRESENÇA\nDIGITAL",
         "Site profissional\nFunil de captação\nConteúdo automatizado\nSEO local e nacional"),
        ("02", "AUTOMAÇÃO\nINTELIGENTE",
         "n8n como orquestrador\nAlertas em tempo real\nRelató­rios automáticos\nIntegração com CRM"),
        ("03", "SAAS DE\nINTELIGÊNCIA",
         "LexRadar · PoliticaIQ\nRegRisk · OpportunityMap\nCrisisMonitor\nReceita recorrente"),
    ]

    for i, (num, title, body) in enumerate(pilars):
        cx = 1.4 + i * 10.7
        cy = 4.0

        # Card background
        card(sl, cx, cy, 10.0, 12.0, fill=MBLACK)

        # Number badge
        rect(sl, cx + 0.35, cy + 0.5, 2.2, 2.2, RED)
        txbox(sl, num,
              x=cx + 0.35, y=cy + 0.55, w=2.2, h=2.2,
              size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        txbox(sl, title,
              x=cx + 0.4, y=cy + 3.2, w=9.2, h=2.4,
              size=19, bold=True, color=WHITE)
        line_sep(sl, cx + 0.4, cy + 5.7, 8.5, RED, thickness=0.08)
        txbox(sl, body,
              x=cx + 0.4, y=cy + 6.0, w=9.2, h=5.5,
              size=13, color=LGRAY)

    return sl


def slide_06_automacoes(prs):
    """Slide 6 – Mapa de Automações"""
    sl = blank_slide(prs)
    section_header_bar(sl, "Automações Prioritárias",
                       "Processos que serão automatizados com n8n + Claude API")
    footer_bar(sl)

    automacoes = [
        ("MONITORAMENTO LEGISLATIVO",
         "Diário • Prioritário",
         "Scraping DOU/Câmara → Classificação IA → Alerta WhatsApp/Email"),
        ("GERAÇÃO DE RELATÓRIOS",
         "Semanal/Mensal • Alto Impacto",
         "Coleta de dados → Análise Claude API → Docx/PDF formatado → Envio automático"),
        ("CAPTURA DE LEADS",
         "Contínuo • Estratégico",
         "Formulário site → CRM → Email sequence → Qualificação IA"),
        ("AGENDAMENTO DE REUNIÕES",
         "Sob demanda • Operacional",
         "Calendly integrado → Confirmação automática → Briefing pré-reunião"),
        ("GESTÃO DE CLIENTES",
         "Contínuo • Recorrente",
         "Painel por cliente → Alertas customizados → NPS automatizado"),
    ]

    for i, (title, tag, flow) in enumerate(automacoes):
        cy = 4.0 + i * 2.55
        card(sl, 1.4, cy, 31.0, 2.2, fill=MBLACK)
        rect(sl, 1.4, cy, 0.45, 2.2, RED)

        txbox(sl, title,
              x=2.3, y=cy + 0.2, w=12, h=1.0,
              size=13, bold=True, color=WHITE)
        txbox(sl, tag,
              x=2.3, y=cy + 1.1, w=12, h=0.8,
              size=10, color=RED)
        txbox(sl, flow,
              x=14.8, y=cy + 0.35, w=17, h=1.5,
              size=11, color=LGRAY)

    return sl


def slide_07_antes_depois(prs):
    """Slide 7 – Antes × Depois"""
    sl = blank_slide(prs)
    section_header_bar(sl, "Transformação Operacional",
                       "Antes × Depois da implantação LocalRise")
    footer_bar(sl)

    # Column headers
    rect(sl, 1.4, 4.0, 14.8, 1.2, DRED)
    txbox(sl, "ANTES — Capitol hoje",
          x=1.4, y=4.05, w=14.8, h=1.1,
          size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rect(sl, 17.0, 4.0, 16.2, 1.2, RED)
    txbox(sl, "DEPOIS — Com LocalRise",
          x=17.0, y=4.05, w=16.2, h=1.1,
          size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    rows = [
        ("Monitoramento manual, 8–10h/semana",
         "Monitoramento automático 24/7, alertas em minutos"),
        ("Relatórios criados do zero a cada entrega",
         "Relatórios gerados automaticamente com 1 clique"),
        ("Sem presença digital, só indicações",
         "Site, funil, SEO + captação ativa de leads"),
        ("Sem receita recorrente",
         "SaaS com assinaturas mensais previsíveis"),
        ("Clientes sem painel de acompanhamento",
         "Dashboard personalizado por cliente, tempo real"),
        ("Sem diferencial tecnológico",
         "Pioneiros em IA regulatória no Brasil"),
    ]

    for i, (before, after) in enumerate(rows):
        cy = 5.5 + i * 1.85
        card(sl, 1.4, cy, 14.8, 1.5, fill=MBLACK)
        txbox(sl, f"✗  {before}",
              x=1.8, y=cy + 0.2, w=14.0, h=1.2,
              size=11, color=LGRAY)
        card(sl, 17.0, cy, 16.2, 1.5, fill=MBLACK, border_color=RED)
        txbox(sl, f"✓  {after}",
              x=17.4, y=cy + 0.2, w=15.5, h=1.2,
              size=11, color=WHITE)

    # VS divider
    rect(sl, 16.1, 4.0, 0.7, 14.0, RED)
    txbox(sl, "VS",
          x=15.8, y=10.5, w=1.4, h=1.2,
          size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    return sl


def slide_08_produtos(prs):
    """Slide 8 – Produtos SaaS"""
    sl = blank_slide(prs)
    section_header_bar(sl, "Produtos SaaS de Inteligência",
                       "Portfolio de produtos desenvolvidos sob a marca Capitol ou white-label LocalRise")
    footer_bar(sl)

    products = [
        ("LexRadar",         "Monitoramento legislativo em tempo real. Alertas configuráveis por tema, órgão e palavra-chave.",        "Assinatura / por cliente"),
        ("PoliticaIQ",       "Análise de posicionamento de parlamentares com IA. Mapas de influência e histórico de votações.",          "Assinatura / relatório"),
        ("RegRisk",          "Score de risco regulatório por setor. Identifica ameaças antes que se tornem problemas.",                 "Consultivo + SaaS"),
        ("OpportunityMap",   "Radar de licitações, PPPs e programas públicos alinhados ao perfil de cada cliente.",                     "Comissão + assinatura"),
        ("CrisisMonitor",    "Monitoramento de mídia e redes sociais para detectar crises regulatórias emergentes.",                    "Retainer mensal"),
    ]

    for i, (name, desc, model) in enumerate(products):
        cy = 4.0 + i * 2.6
        card(sl, 1.4, cy, 31.0, 2.3, fill=MBLACK)
        rect(sl, 1.4, cy, 0.45, 2.3, RED)

        txbox(sl, name,
              x=2.3, y=cy + 0.2, w=9, h=1.0,
              size=15, bold=True, color=RED)
        txbox(sl, desc,
              x=2.3, y=cy + 1.1, w=20, h=1.0,
              size=11, color=LGRAY)
        txbox(sl, model,
              x=23.5, y=cy + 0.5, w=9, h=1.2,
              size=11, bold=True, color=WHITE, align=PP_ALIGN.RIGHT)

    return sl


def slide_09_pipeline(prs):
    """Slide 9 – Pipeline de Inteligência: 6 Agentes"""
    sl = blank_slide(prs)
    section_header_bar(sl, "Pipeline de Inteligência Regulatória",
                       "6 agentes em série — da coleta ao cliente final")
    footer_bar(sl)

    agents = [
        ("01", "COLLECTOR",   "Scraping DOU,\nCâmara, Senado,\nTCU e portais"),
        ("02", "CLASSIFIER",  "Classifica por\ntema, urgência\ne relevância"),
        ("03", "ANALYST",     "Claude API gera\nanálise técnica\ne contexto"),
        ("04", "RISK SCORER", "Score 0–100\npor setor e\nperfil do cliente"),
        ("05", "REPORTER",    "Relatório Word/\nPDF formatado\nautomaticamente"),
        ("06", "DISPATCHER",  "Entrega via\nEmail, WhatsApp\nou dashboard"),
    ]

    card_w = 4.8
    gap = 0.6
    start_x = 1.0
    cy = 4.5

    for i, (num, name, desc) in enumerate(agents):
        cx = start_x + i * (card_w + gap)
        card(sl, cx, cy, card_w, 9.5, fill=MBLACK)
        rect(sl, cx, cy, card_w, 0.4, RED)

        txbox(sl, num,
              x=cx + 0.3, y=cy + 0.7, w=card_w - 0.6, h=1.5,
              size=26, bold=True, color=DRED, align=PP_ALIGN.CENTER)
        txbox(sl, name,
              x=cx + 0.2, y=cy + 2.3, w=card_w - 0.4, h=1.5,
              size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        line_sep(sl, cx + 0.5, cy + 3.9, card_w - 1.0, RED, thickness=0.07)
        txbox(sl, desc,
              x=cx + 0.3, y=cy + 4.2, w=card_w - 0.6, h=4.5,
              size=12, color=LGRAY, align=PP_ALIGN.CENTER)

        # Arrow connector (except last)
        if i < len(agents) - 1:
            ax = cx + card_w + 0.05
            txbox(sl, "→",
                  x=ax, y=cy + 4.0, w=gap, h=1.5,
                  size=18, bold=True, color=RED, align=PP_ALIGN.CENTER)

    # n8n label
    rect(sl, 1.0, 14.6, 31.4, 0.9, DRED)
    txbox(sl, "Orquestrador: n8n  ·  LLM: Claude API (Anthropic)  ·  Entrega: Multi-canal",
          x=1.5, y=14.7, w=30, h=0.7,
          size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    return sl


def slide_10_arquitetura(prs):
    """Slide 10 – Stack Tecnológico"""
    sl = blank_slide(prs)
    section_header_bar(sl, "Arquitetura Tecnológica",
                       "Stack moderno, escalável e de baixo custo operacional")
    footer_bar(sl)

    layers = [
        ("CAMADA DE COLETA",    "Scrapers Python · RSS · APIs públicas (Câmara, Senado, DOU) · Web crawlers"),
        ("CAMADA DE ORQUESTRAÇÃO", "n8n self-hosted — Workflows visuais, agendamentos, integrações e lógica de negócio"),
        ("CAMADA DE INTELIGÊNCIA", "Claude API (Anthropic) — Análise, classificação, scoring, geração de texto e relatórios"),
        ("CAMADA DE DADOS",     "Supabase (PostgreSQL) — Armazenamento, histórico, perfis de clientes e logs"),
        ("CAMADA DE ENTREGA",   "Email (SMTP) · WhatsApp (Evolution API) · Dashboard Next.js · PDF/Docx gerado"),
    ]

    for i, (layer, desc) in enumerate(layers):
        cy = 4.2 + i * 2.55
        fill = MBLACK if i % 2 == 0 else RGBColor(0x16, 0x16, 0x16)
        card(sl, 1.4, cy, 31.0, 2.2, fill=fill)
        rect(sl, 1.4, cy, 0.45, 2.2, RED)
        txbox(sl, layer,
              x=2.3, y=cy + 0.2, w=10, h=1.0,
              size=12, bold=True, color=RED)
        txbox(sl, desc,
              x=13.0, y=cy + 0.3, w=19.5, h=1.6,
              size=12, color=LGRAY)

    return sl


def slide_11_fases(prs):
    """Slide 11 – Plano de Implementação"""
    sl = blank_slide(prs)
    section_header_bar(sl, "Plano de Implementação",
                       "Quatro fases progressivas — do diagnóstico ao produto completo")
    footer_bar(sl)

    phases = [
        ("FASE 1", "Fundação Digital", "30 dias",
         "Site profissional\nSetup Google Analytics\nPerfis redes sociais\nIdentidade visual digital"),
        ("FASE 2", "Automação Core", "60 dias",
         "Monitoramento legislativo automatizado\nAlertas configurados por cliente\nPrimeiros relatórios automáticos\nCRM integrado"),
        ("FASE 3", "SaaS MVP", "90 dias",
         "LexRadar em produção\nDashboard por cliente\nOnboarding dos primeiros assinantes\nDocumentação e suporte"),
        ("FASE 4", "Escala", "180 dias",
         "Portfolio completo (5 produtos)\nTime de CS\nParcerias estratégicas\nReceita recorrente consolidada"),
    ]

    for i, (phase, title, deadline, items) in enumerate(phases):
        col = i % 2
        row = i // 2
        cx = 1.4 + col * 16.0
        cy = 4.2 + row * 6.5

        card(sl, cx, cy, 15.0, 6.0, fill=MBLACK)
        rect(sl, cx, cy, 15.0, 0.45, RED)

        txbox(sl, phase,
              x=cx + 0.4, y=cy + 0.7, w=6, h=1.1,
              size=22, bold=True, color=RED)
        txbox(sl, deadline,
              x=cx + 10, y=cy + 0.7, w=4.5, h=1.1,
              size=14, bold=True, color=GRAY, align=PP_ALIGN.RIGHT)
        txbox(sl, title,
              x=cx + 0.4, y=cy + 1.8, w=14, h=1.2,
              size=16, bold=True, color=WHITE)
        txbox(sl, items,
              x=cx + 0.4, y=cy + 3.0, w=14, h=2.8,
              size=12, color=LGRAY)

    return sl


def slide_12_roi(prs):
    """Slide 12 – ROI e Métricas Esperadas"""
    sl = blank_slide(prs)
    section_header_bar(sl, "ROI & Métricas Esperadas",
                       "Projeções conservadoras para os primeiros 12 meses")
    footer_bar(sl)

    kpis = [
        ("40h/mês",   "Horas economizadas\nem monitoramento manual"),
        ("3–5×",      "Ampliação do alcance\nde prospecção digital"),
        ("R$ 8–15k",  "Novos contratos mensais\nvia inbound marketing"),
        ("R$ 5–20k",  "MRR potencial\nde assinaturas SaaS"),
        ("< 90 dias", "Payback estimado\ndo investimento"),
        ("10–30 clientes", "Base assinante\nno primeiro ano"),
    ]

    for i, (kpi, desc) in enumerate(kpis):
        col = i % 3
        row = i // 3
        cx = 1.4 + col * 10.7
        cy = 4.5 + row * 5.5

        card(sl, cx, cy, 9.8, 4.8, fill=MBLACK, border_color=RED)
        txbox(sl, kpi,
              x=cx + 0.4, y=cy + 0.8, w=9, h=2.2,
              size=30, bold=True, color=RED, align=PP_ALIGN.CENTER)
        line_sep(sl, cx + 1.0, cy + 3.0, 7.5, DRED, thickness=0.08)
        txbox(sl, desc,
              x=cx + 0.4, y=cy + 3.2, w=9, h=1.4,
              size=12, color=LGRAY, align=PP_ALIGN.CENTER)

    return sl


def slide_13_modelo_comercial(prs):
    """Slide 13 – Modelo Comercial"""
    sl = blank_slide(prs)
    section_header_bar(sl, "Modelo Comercial",
                       "Três modalidades de engajamento com a LocalRise")
    footer_bar(sl)

    plans = [
        ("STARTER",   "Implantação",
         "Site + identidade digital\nSetup automação básica\nPrimeiro relatório automático",
         "A partir de\nR$ 3.500", False),
        ("GROWTH",    "Recomendado",
         "Tudo do Starter\n+ LexRadar MVP\n+ Dashboard de clientes\n+ 3 meses de suporte",
         "A partir de\nR$ 6.800", True),
        ("ENTERPRISE","Parceria Plena",
         "Tudo do Growth\n+ Portfolio completo\n+ White-label\n+ Revenue share SaaS",
         "Sob consulta", False),
    ]

    for i, (name, tag, items, price, highlight) in enumerate(plans):
        cx = 1.4 + i * 10.7
        cy = 4.2

        fill = RED if highlight else MBLACK
        card(sl, cx, cy, 10.0, 12.5, fill=fill)

        if highlight:
            rect(sl, cx, cy, 10.0, 0.5, DRED)
            txbox(sl, "★ MAIS ESCOLHIDO",
                  x=cx + 0.3, y=cy + 0.55, w=9.4, h=0.9,
                  size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        else:
            rect(sl, cx, cy, 10.0, 0.5, DRED)

        txbox(sl, name,
              x=cx + 0.4, y=cy + 1.6, w=9.2, h=1.5,
              size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(sl, tag,
              x=cx + 0.4, y=cy + 3.0, w=9.2, h=0.9,
              size=12, color=LGRAY if not highlight else WHITE, align=PP_ALIGN.CENTER)
        line_sep(sl, cx + 1.0, cy + 4.0, 8.0, WHITE if highlight else RED, thickness=0.07)
        txbox(sl, items,
              x=cx + 0.5, y=cy + 4.3, w=9.0, h=5.5,
              size=12, color=WHITE)
        txbox(sl, price,
              x=cx + 0.4, y=cy + 10.0, w=9.2, h=2.0,
              size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    return sl


def slide_14_por_que_localrise(prs):
    """Slide 14 – Por que LocalRise"""
    sl = blank_slide(prs)
    section_header_bar(sl, "Por que LocalRise Advisory?",
                       "O parceiro certo para esta transformação")
    footer_bar(sl)

    reasons = [
        ("Especialistas em Gov-Tech",
         "Profundo conhecimento do mercado de relações governamentais, compliance e inteligência regulatória."),
        ("Stack proprietária",
         "Automações desenvolvidas internamente com n8n + Claude API — sem dependência de ferramentas de terceiros caras."),
        ("Entrega orientada a resultado",
         "Não vendemos horas. Vendemos outcomes: leads, relatórios, alertas e receita recorrente gerados."),
        ("Modelo de parceria",
         "Alinhamento via revenue share em produtos SaaS — nosso sucesso é o seu sucesso."),
        ("Velocidade de implantação",
         "MVP rodando em 30–60 dias. Sem projetos longos de TI. Sem burocracia."),
    ]

    for i, (title, body) in enumerate(reasons):
        cy = 4.2 + i * 2.6
        card(sl, 1.4, cy, 31.0, 2.2, fill=MBLACK)
        rect(sl, 1.4, cy, 0.45, 2.2, RED)

        txbox(sl, f"0{i+1}  {title}",
              x=2.3, y=cy + 0.2, w=13, h=1.0,
              size=14, bold=True, color=RED)
        txbox(sl, body,
              x=2.3, y=cy + 1.1, w=29.5, h=1.0,
              size=12, color=LGRAY)

    return sl


def slide_15_cta(prs):
    """Slide 15 – CTA Final"""
    sl = blank_slide(prs)

    # Split background
    rect(sl, 0, 0, 16.5, 19.05, BLACK)
    rect(sl, 16.5, 0, 17.367, 19.05, RED)
    rect(sl, 16.0, 0, 0.8, 19.05, DRED)

    # Left side
    add_logo(sl, x=1.5, y=1.0, w=4.0)
    txbox(sl, "Próximos\nPassos",
          x=1.5, y=3.2, w=14, h=4.5,
          size=44, bold=True, color=WHITE)
    line_sep(sl, 1.5, 7.8, 13, RED, thickness=0.12)

    steps = [
        "1.  Validar escopo e prioridades (esta reunião)",
        "2.  Proposta formal com cronograma detalhado",
        "3.  Kickoff em até 15 dias após assinatura",
        "4.  Primeiro entregável em até 30 dias",
    ]
    txbox(sl, "\n".join(steps),
          x=1.5, y=8.2, w=14, h=6.5,
          size=13, color=LGRAY)

    txbox(sl, f"Sao Paulo, {DATE}",
          x=1.5, y=15.5, w=14, h=1.0,
          size=11, color=GRAY)

    # Right side
    txbox(sl, "Vamos\nconstruir\njuntos?",
          x=17.2, y=2.5, w=15, h=6.5,
          size=40, bold=True, color=WHITE)

    txbox(sl, "LocalRise Advisory",
          x=17.2, y=10.0, w=15, h=1.4,
          size=20, bold=True, color=WHITE)
    txbox(sl, "Nicolas Slater\nDiretoria Comercial",
          x=17.2, y=11.4, w=15, h=2.0,
          size=14, color=WHITE)
    line_sep(sl, 17.2, 13.7, 14, WHITE, thickness=0.08)
    txbox(sl, "contato@localrise.com.br\nlocalrise.com.br",
          x=17.2, y=14.0, w=15, h=2.0,
          size=14, color=WHITE)

    return sl


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()

    print("Gerando slides...")
    slide_01_cover(prs)            ; print("  OK Slide 01 - Capa")
    slide_02_sumario(prs)          ; print("  OK Slide 02 - Sumario Executivo")
    slide_03_diagnostico(prs)      ; print("  OK Slide 03 - Diagnostico")
    slide_04_oportunidade(prs)     ; print("  OK Slide 04 - Oportunidade")
    slide_05_solucao(prs)          ; print("  OK Slide 05 - Solucao")
    slide_06_automacoes(prs)       ; print("  OK Slide 06 - Automacoes")
    slide_07_antes_depois(prs)     ; print("  OK Slide 07 - Antes x Depois")
    slide_08_produtos(prs)         ; print("  OK Slide 08 - Produtos SaaS")
    slide_09_pipeline(prs)         ; print("  OK Slide 09 - Pipeline de IA")
    slide_10_arquitetura(prs)      ; print("  OK Slide 10 - Arquitetura")
    slide_11_fases(prs)            ; print("  OK Slide 11 - Fases")
    slide_12_roi(prs)              ; print("  OK Slide 12 - ROI")
    slide_13_modelo_comercial(prs) ; print("  OK Slide 13 - Modelo Comercial")
    slide_14_por_que_localrise(prs); print("  OK Slide 14 - Por que LocalRise")
    slide_15_cta(prs)              ; print("  OK Slide 15 - CTA Final")

    prs.save(str(OUT_PATH))
    print("\nPronto! Apresentacao salva em:")
    print("  " + str(OUT_PATH))


if __name__ == "__main__":
    build()
